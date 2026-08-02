import os
import re
import fitz
import pptx
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

# 1. Parse PDF: TalkFile_지구촌교회_중등부_OX퀴즈_자료집.pdf
script_dir = os.path.dirname(os.path.abspath(__file__))
pdf_path = os.path.join(script_dir, '2026SummerFestival', 'TalkFile_지구촌교회_중등부_OX퀴즈_자료집.pdf')
if not os.path.exists(pdf_path):
    pdf_path = os.path.join(script_dir, '지구촌교회_중등부_OX퀴즈_자료집.pdf')

doc = fitz.open(pdf_path)
all_qs = {}

for p_idx in range(len(doc)):
    page = doc[p_idx]
    y_lines = []
    for d in page.get_drawings():
        r = d.get('rect')
        if r and (r.y1 - r.y0) <= 2.5 and r.y0 >= 30 and r.y1 <= 805:
            if abs(r.x0 - 42.8) < 5 or abs(r.x0 - 43.5) < 5:
                y_lines.append(round(r.y0, 1))

    y_lines = sorted(list(set(y_lines)))
    row_ys = [y for y in y_lines if y >= 50]
    row_ys.append(795.0)
    
    intervals = []
    for i in range(len(row_ys) - 1):
        if row_ys[i+1] - row_ys[i] >= 10:
            intervals.append((row_ys[i], row_ys[i+1]))

    for y_top, y_bot in intervals:
        def get_cell(x_start, x_end):
            c_rect = fitz.Rect(x_start, y_top + 1, x_end, y_bot - 1)
            cell_words = page.get_text('words', clip=c_rect)
            lines = {}
            for w in cell_words:
                yc = round((w[1] + w[3]) / 2, 1)
                matched_y = None
                for ly in lines:
                    if abs(ly - yc) < 3.0:
                        matched_y = ly
                        break
                if matched_y is None:
                    matched_y = yc
                    lines[matched_y] = []
                lines[matched_y].append(w)
                
            res_lines = []
            for ly in sorted(lines.keys()):
                lw = sorted(lines[ly], key=lambda w: w[0])
                res_lines.append(' '.join([w[4] for w in lw]))
            return ' '.join(res_lines).strip()

        no_str = get_cell(42.8, 69.0)
        if no_str.isdigit():
            num = int(no_str)
            diff = get_cell(69.0, 105.0)
            cat = get_cell(105.0, 160.5)
            q_text = get_cell(160.5, 385.5)
            ans = get_cell(385.5, 420.8)
            exp = get_cell(420.8, 553.5)
            
            all_qs[num] = {
                'pdf_no': num,
                'difficulty': diff,
                'category': cat,
                'question': q_text,
                'answer': ans,
                'explanation': exp
            }

main_seq = [6, 17, 21, 7, 31, 13, 25, 29, 10, 37, 20, 39, 15, 43, 63, 47, 32, 65, 49, 61, 51, 67, 58, 71, 68, 74, 84, 77, 86, 79, 93, 89, 94, 91, 83]
extra_seq = [98, 100, 81, 82, 85, 18, 22, 26, 53, 95]

main_items = [all_qs[n] for n in main_seq]
extra_items = [all_qs[n] for n in extra_seq]

print(f'Loaded {len(main_items)} main items and {len(extra_items)} extra items from TalkFile PDF.')

# 2. Build PPTX
prs = pptx.Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

blank_layout = prs.slide_layouts[6]

# Theme Colors
C_DARK_BG = RGBColor(27, 54, 93)      # #1B365D Navy
C_LIGHT_BG = RGBColor(248, 250, 252)  # #F8FAFC Off-white
C_CARD_BG = RGBColor(255, 255, 255)   # #FFFFFF White
C_NAVY_TEXT = RGBColor(27, 54, 93)
C_DARK_TEXT = RGBColor(30, 41, 59)
C_MUTED_TEXT = RGBColor(100, 116, 139)
C_GOLD = RGBColor(245, 158, 11)        # #F59E0B Gold
C_BLUE_BADGE = RGBColor(43, 84, 143)  # #2B548F Royal Blue

C_O_COLOR = RGBColor(16, 185, 129)     # #10B981 Green
C_X_COLOR = RGBColor(239, 68, 68)      # #EF4444 Red

FONT_FAMILY = 'Malgun Gothic'

def add_solid_bg(slide, color):
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(13.333), Inches(7.5))
    bg.fill.solid()
    bg.fill.fore_color.rgb = color
    bg.line.fill.background()
    return bg

# Slide 1: Title Slide
slide = prs.slides.add_slide(blank_layout)
add_solid_bg(slide, C_DARK_BG)

tx_box = slide.shapes.add_textbox(Inches(1.0), Inches(2.0), Inches(11.333), Inches(3.5))
tf = tx_box.text_frame
tf.word_wrap = True

p0 = tf.paragraphs[0]
p0.text = '지구촌교회 청소년지구 중등부'
p0.font.name = FONT_FAMILY
p0.font.size = Pt(28)
p0.font.bold = True
p0.font.color.rgb = C_GOLD
p0.alignment = PP_ALIGN.CENTER

p1 = tf.add_paragraph()
p1.text = 'O / X  퀴 즈  대 회'
p1.font.name = FONT_FAMILY
p1.font.size = Pt(56)
p1.font.bold = True
p1.font.color.rgb = RGBColor(255, 255, 255)
p1.alignment = PP_ALIGN.CENTER
p1.space_before = Pt(14)

p2 = tf.add_paragraph()
p2.text = '수련회 & 진행용 맞춤 퀴즈 PPT  |  총 45문항 (본 퀴즈 35문항 + 여분 10문항)'
p2.font.name = FONT_FAMILY
p2.font.size = Pt(20)
p2.font.color.rgb = RGBColor(203, 213, 225)
p2.alignment = PP_ALIGN.CENTER
p2.space_before = Pt(24)

# Section Slide helper
def add_section_slide(title, subtitle, section_num):
    s = prs.slides.add_slide(blank_layout)
    add_solid_bg(s, C_DARK_BG)
    
    tb = s.shapes.add_textbox(Inches(1.0), Inches(2.5), Inches(11.333), Inches(3.0))
    tf = tb.text_frame
    tf.word_wrap = True
    
    p = tf.paragraphs[0]
    p.text = f'PART {section_num}'
    p.font.name = FONT_FAMILY
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = C_GOLD
    p.alignment = PP_ALIGN.CENTER
    
    p2 = tf.add_paragraph()
    p2.text = title
    p2.font.name = FONT_FAMILY
    p2.font.size = Pt(48)
    p2.font.bold = True
    p2.font.color.rgb = RGBColor(255, 255, 255)
    p2.alignment = PP_ALIGN.CENTER
    p2.space_before = Pt(10)
    
    p3 = tf.add_paragraph()
    p3.text = subtitle
    p3.font.name = FONT_FAMILY
    p3.font.size = Pt(20)
    p3.font.color.rgb = RGBColor(203, 213, 225)
    p3.alignment = PP_ALIGN.CENTER
    p3.space_before = Pt(16)

add_section_slide('본 퀴즈 (1번 ~ 35번)', '준비되셨나요? 문제에 집중하고 O / X 판을 올려주세요!', 1)

def create_badge(slide, left, top, width, height, text, bg_color, text_color=RGBColor(255,255,255), font_size=15):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = bg_color
    shape.line.fill.background()
    tf = shape.text_frame
    tf.word_wrap = False
    p = tf.paragraphs[0]
    p.text = text
    p.font.name = FONT_FAMILY
    p.font.size = Pt(font_size)
    p.font.bold = True
    p.font.color.rgb = text_color
    p.alignment = PP_ALIGN.CENTER
    return shape

def build_question_slides(item, q_idx, total_count, is_extra=False):
    prefix = '여분 Q' if is_extra else 'Q'
    q_badge_text = f'{prefix} {q_idx:02d} / {total_count:02d}'
    
    # 1) QUESTION SLIDE
    s_q = prs.slides.add_slide(blank_layout)
    add_solid_bg(s_q, C_LIGHT_BG)
    
    top_bar = s_q.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(13.333), Inches(1.1))
    top_bar.fill.solid()
    top_bar.fill.fore_color.rgb = C_DARK_BG
    top_bar.line.fill.background()
    
    create_badge(s_q, Inches(0.8), Inches(0.28), Inches(2.2), Inches(0.54), q_badge_text, C_GOLD, C_DARK_BG, 17)
    create_badge(s_q, Inches(3.2), Inches(0.28), Inches(2.0), Inches(0.54), f'[ {item["category"]} ]', C_BLUE_BADGE, RGBColor(255,255,255), 16)
    create_badge(s_q, Inches(5.4), Inches(0.28), Inches(1.8), Inches(0.54), f'난이도: {item["difficulty"]}', RGBColor(71, 85, 105), RGBColor(255,255,255), 15)
    create_badge(s_q, Inches(10.5), Inches(0.28), Inches(2.0), Inches(0.54), f'자료집 #{item["pdf_no"]}', RGBColor(51, 65, 85), RGBColor(226, 232, 240), 14)
    
    card = s_q.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(1.5), Inches(11.733), Inches(3.6))
    card.fill.solid()
    card.fill.fore_color.rgb = C_CARD_BG
    card.line.color.rgb = RGBColor(226, 232, 240)
    card.line.width = Pt(1.5)
    
    q_box = s_q.shapes.add_textbox(Inches(1.2), Inches(1.7), Inches(10.933), Inches(3.2))
    tf_q = q_box.text_frame
    tf_q.word_wrap = True
    
    pq = tf_q.paragraphs[0]
    pq.text = item['question']
    pq.font.name = FONT_FAMILY
    q_len = len(item['question'])
    if q_len > 70:
        pq.font.size = Pt(28)
    elif q_len > 45:
        pq.font.size = Pt(32)
    else:
        pq.font.size = Pt(36)
    pq.font.bold = True
    pq.font.color.rgb = C_DARK_TEXT
    pq.alignment = PP_ALIGN.CENTER
    
    box_o = s_q.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(2.2), Inches(5.4), Inches(4.0), Inches(1.5))
    box_o.fill.solid()
    box_o.fill.fore_color.rgb = RGBColor(236, 253, 245)
    box_o.line.color.rgb = C_O_COLOR
    box_o.line.width = Pt(3)
    
    tf_o = box_o.text_frame
    p_o = tf_o.paragraphs[0]
    p_o.text = 'O'
    p_o.font.name = FONT_FAMILY
    p_o.font.size = Pt(64)
    p_o.font.bold = True
    p_o.font.color.rgb = C_O_COLOR
    p_o.alignment = PP_ALIGN.CENTER
    
    box_x = s_q.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(7.133), Inches(5.4), Inches(4.0), Inches(1.5))
    box_x.fill.solid()
    box_x.fill.fore_color.rgb = RGBColor(254, 242, 242)
    box_x.line.color.rgb = C_X_COLOR
    box_x.line.width = Pt(3)
    
    tf_x = box_x.text_frame
    p_x = tf_x.paragraphs[0]
    p_x.text = 'X'
    p_x.font.name = FONT_FAMILY
    p_x.font.size = Pt(64)
    p_x.font.bold = True
    p_x.font.color.rgb = C_X_COLOR
    p_x.alignment = PP_ALIGN.CENTER

    # 2) ANSWER SLIDE
    s_a = prs.slides.add_slide(blank_layout)
    add_solid_bg(s_a, C_LIGHT_BG)
    
    top_bar_a = s_a.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(13.333), Inches(1.1))
    top_bar_a.fill.solid()
    top_bar_a.fill.fore_color.rgb = C_DARK_BG
    top_bar_a.line.fill.background()
    
    create_badge(s_a, Inches(0.8), Inches(0.28), Inches(2.2), Inches(0.54), q_badge_text, C_GOLD, C_DARK_BG, 17)
    create_badge(s_a, Inches(3.2), Inches(0.28), Inches(2.0), Inches(0.54), f'[ {item["category"]} ]', C_BLUE_BADGE, RGBColor(255,255,255), 16)
    create_badge(s_a, Inches(5.4), Inches(0.28), Inches(2.2), Inches(0.54), '정답 및 해설', RGBColor(16, 185, 129), RGBColor(255,255,255), 16)
    create_badge(s_a, Inches(10.5), Inches(0.28), Inches(2.0), Inches(0.54), f'자료집 #{item["pdf_no"]}', RGBColor(51, 65, 85), RGBColor(226, 232, 240), 14)
    
    q_sum_card = s_a.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(1.3), Inches(11.733), Inches(1.0))
    q_sum_card.fill.solid()
    q_sum_card.fill.fore_color.rgb = C_CARD_BG
    q_sum_card.line.color.rgb = RGBColor(226, 232, 240)
    
    tf_sum = q_sum_card.text_frame
    tf_sum.word_wrap = True
    p_sum = tf_sum.paragraphs[0]
    p_sum.text = f'Q. {item["question"]}'
    p_sum.font.name = FONT_FAMILY
    p_sum.font.size = Pt(18)
    p_sum.font.bold = True
    p_sum.font.color.rgb = C_MUTED_TEXT
    p_sum.alignment = PP_ALIGN.LEFT
    
    ans_is_o = (item['answer'].upper() == 'O')
    ans_color = C_O_COLOR if ans_is_o else C_X_COLOR
    ans_bg = RGBColor(236, 253, 245) if ans_is_o else RGBColor(254, 242, 242)
    
    ans_box = s_a.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(2.5), Inches(3.8), Inches(4.5))
    ans_box.fill.solid()
    ans_box.fill.fore_color.rgb = ans_bg
    ans_box.line.color.rgb = ans_color
    ans_box.line.width = Pt(3)
    
    tf_ans = ans_box.text_frame
    tf_ans.word_wrap = True
    
    pa1 = tf_ans.paragraphs[0]
    pa1.text = '정 답'
    pa1.font.name = FONT_FAMILY
    pa1.font.size = Pt(22)
    pa1.font.bold = True
    pa1.font.color.rgb = C_DARK_TEXT
    pa1.alignment = PP_ALIGN.CENTER
    
    pa2 = tf_ans.add_paragraph()
    pa2.text = item['answer']
    pa2.font.name = FONT_FAMILY
    pa2.font.size = Pt(120)
    pa2.font.bold = True
    pa2.font.color.rgb = ans_color
    pa2.alignment = PP_ALIGN.CENTER
    
    exp_card = s_a.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(4.9), Inches(2.5), Inches(7.633), Inches(4.5))
    exp_card.fill.solid()
    exp_card.fill.fore_color.rgb = C_CARD_BG
    exp_card.line.color.rgb = RGBColor(226, 232, 240)
    exp_card.line.width = Pt(1.5)
    
    exp_box = s_a.shapes.add_textbox(Inches(5.2), Inches(2.7), Inches(7.033), Inches(4.1))
    tf_exp = exp_box.text_frame
    tf_exp.word_wrap = True
    
    pe_title = tf_exp.paragraphs[0]
    pe_title.text = '💡 해설 및 참고'
    pe_title.font.name = FONT_FAMILY
    pe_title.font.size = Pt(24)
    pe_title.font.bold = True
    pe_title.font.color.rgb = C_NAVY_TEXT
    pe_title.space_after = Pt(14)
    
    pe_text = tf_exp.add_paragraph()
    pe_text.text = item['explanation']
    pe_text.font.name = FONT_FAMILY
    
    exp_len = len(item['explanation'])
    if exp_len > 80:
        pe_text.font.size = Pt(22)
    elif exp_len > 40:
        pe_text.font.size = Pt(26)
    else:
        pe_text.font.size = Pt(30)
    pe_text.font.bold = True
    pe_text.font.color.rgb = C_DARK_TEXT

for idx, item in enumerate(main_items, 1):
    build_question_slides(item, idx, len(main_items), is_extra=False)

add_section_slide('여분 문제 (1번 ~ 10번)', '패자부활전 & 동점자 처리용 예비 퀴즈 10문항', 2)

for idx, item in enumerate(extra_items, 1):
    build_question_slides(item, idx, len(extra_items), is_extra=True)

slide_end = prs.slides.add_slide(blank_layout)
add_solid_bg(slide_end, C_DARK_BG)

tx_box_end = slide_end.shapes.add_textbox(Inches(1.0), Inches(2.5), Inches(11.333), Inches(3.0))
tf_end = tx_box_end.text_frame
tf_end.word_wrap = True

pe1 = tf_end.paragraphs[0]
pe1.text = '수고하셨습니다!'
pe1.font.name = FONT_FAMILY
pe1.font.size = Pt(56)
pe1.font.bold = True
pe1.font.color.rgb = C_GOLD
pe1.alignment = PP_ALIGN.CENTER

pe2 = tf_end.add_paragraph()
pe2.text = '지구촌교회 청소년지구 중등부 O/X 퀴즈 종료'
pe2.font.name = FONT_FAMILY
pe2.font.size = Pt(26)
pe2.font.color.rgb = RGBColor(255, 255, 255)
pe2.alignment = PP_ALIGN.CENTER
pe2.space_before = Pt(20)

target_pptx = os.path.join(script_dir, '지구촌교회_중등부_OX퀴즈_진행용.pptx')
prs.save(target_pptx)
print(f'Successfully updated presentation: {target_pptx}')
print(f'Total slides: {len(prs.slides)}')
